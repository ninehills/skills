#!/usr/bin/env python3
"""Inject manuscript content into PPTX speaker notes.

Reads per-page manuscript from output.json and appends it to each slide's
speaker notes as <manuscript>...</manuscript>, preserving existing notes.

Usage:
    python3 inject_notes.py <input_pptx> <output_json> <output_pptx>

Dependencies: python-pptx (already required by the skill)
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


MANUSCRIPT_OPEN = "<manuscript>"
MANUSCRIPT_CLOSE = "</manuscript>"


def load_page_content(json_path: str) -> dict[int, str]:
    """Load output.json and build a page_number -> content mapping."""
    with open(json_path, encoding="utf-8") as f:
        data = json.load(f)

    page_map = {}
    for section in data.get("sections", []):
        for page in section.get("pages", []):
            page_num = page.get("page")
            content = page.get("content", "").strip()
            if page_num and content:
                page_map[page_num] = content
    return page_map


def has_existing_manuscript(notes_text: str) -> bool:
    """Check if notes already contain a <manuscript> block."""
    return MANUSCRIPT_OPEN in notes_text


def inject_notes(input_pptx: str, page_map: dict[int, str], output_pptx: str) -> dict:
    """Inject manuscript into speaker notes, return stats."""
    prs = Presentation(input_pptx)
    stats = {"total_slides": len(prs.slides), "injected": 0, "skipped_no_content": 0, "skipped_existing": 0}

    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        content = page_map.get(page_num)

        if not content:
            stats["skipped_no_content"] += 1
            continue

        # Access or create notes slide
        # python-pptx creates the notes slide automatically if it doesn't exist
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame

        # Check existing notes for duplicate manuscript block
        existing_text = tf.text
        if has_existing_manuscript(existing_text):
            stats["skipped_existing"] += 1
            continue

        # Append manuscript block as new paragraphs
        # Add a blank separator line first
        sep = tf.add_paragraph()
        sep.text = ""

        # Add the manuscript wrapped in tags
        manuscript_text = f"{MANUSCRIPT_OPEN}\n{content}\n{MANUSCRIPT_CLOSE}"

        for line in manuscript_text.split("\n"):
            p = tf.add_paragraph()
            p.text = line
            # Keep font small and subtle so it doesn't dominate the notes view
            for run in p.runs:
                run.font.size = Pt(10)

        stats["injected"] += 1

    # Write to a temp file first then move (avoids partial writes)
    import tempfile, shutil
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp.close()
    prs.save(tmp.name)
    shutil.move(tmp.name, output_pptx)

    return stats


def main():
    if len(sys.argv) < 4:
        print("Usage: inject_notes.py <input_pptx> <output_json> <output_pptx>")
        print()
        print("  input_pptx   Original PPTX file (from 00_input/)")
        print("  output_json  Manuscript JSON (from 04_final/output.json)")
        print("  output_pptx  Where to save the annotated PPTX")
        sys.exit(1)

    input_pptx = sys.argv[1]
    output_json = sys.argv[2]
    output_pptx = sys.argv[3]

    # Validate inputs
    if not Path(input_pptx).exists():
        print(f"ERROR: 输入文件不存在: {input_pptx}")
        sys.exit(1)
    if not Path(output_json).exists():
        print(f"ERROR: JSON 文件不存在: {output_json}")
        sys.exit(1)

    # Load manuscript content
    page_map = load_page_content(output_json)
    if not page_map:
        print("WARNING: output.json 中未找到任何页面内容")

    # Inject
    stats = inject_notes(input_pptx, page_map, output_pptx)

    # Report
    print(f"\n{'='*50}")
    print(f"演讲者备注注入完成")
    print(f"{'='*50}")
    print(f"  总页数:       {stats['total_slides']}")
    print(f"  已注入:       {stats['injected']}")
    print(f"  无对应内容:   {stats['skipped_no_content']}")
    print(f"  已有 manuscript: {stats['skipped_existing']}")
    print(f"  输出文件:     {output_pptx}")

    # JSON output
    print("\n--- JSON ---")
    print(json.dumps(stats, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
