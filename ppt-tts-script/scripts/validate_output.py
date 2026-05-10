#!/usr/bin/env python3
"""Validate output of each pipeline step.

Usage:
    python3 validate_output.py <output_dir> --step <1|2|3|4>

Returns exit code 0 if valid, 1 if errors found.
Prints structured report to stdout.
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path


def validate_step1(output_dir: str) -> list[dict]:
    """Validate material extraction output."""
    errors = []
    materials_dir = Path(output_dir) / "01_materials"

    if not materials_dir.exists():
        return [{"severity": "error", "message": "01_materials/ 目录不存在"}]

    # Check index.json
    index_path = materials_dir / "index.json"
    if not index_path.exists():
        errors.append({"severity": "error", "message": "index.json 不存在"})
    else:
        try:
            with open(index_path) as f:
                index = json.load(f)
            total = index.get("total_pages", 0)
            if total == 0:
                errors.append({"severity": "error", "message": "index.json 中 total_pages 为 0"})
        except json.JSONDecodeError as e:
            errors.append({"severity": "error", "message": f"index.json JSON 解析失败: {e}"})
            total = 0

    # Count PNG and MD files
    png_files = sorted(materials_dir.glob("slide-*.png"))
    md_files = sorted(materials_dir.glob("slide-*.md"))

    if len(png_files) == 0:
        errors.append({"severity": "error", "message": "未找到 PNG 文件"})
    if len(md_files) == 0:
        errors.append({"severity": "error", "message": "未找到 MD 文件"})

    if len(png_files) != len(md_files):
        errors.append({
            "severity": "warning",
            "message": f"PNG 文件数 ({len(png_files)}) ≠ MD 文件数 ({len(md_files)})"
        })

    # Check numbering consistency
    png_nums = {int(re.search(r"slide-(\d+)", f.name).group(1)) for f in png_files if re.search(r"slide-(\d+)", f.name)}
    md_nums = {int(re.search(r"slide-(\d+)", f.name).group(1)) for f in md_files if re.search(r"slide-(\d+)", f.name)}

    missing_png = md_nums - png_nums
    missing_md = png_nums - md_nums
    if missing_png:
        errors.append({"severity": "warning", "message": f"以下页码有 MD 但无 PNG: {sorted(missing_png)}"})
    if missing_md:
        errors.append({"severity": "warning", "message": f"以下页码有 PNG 但无 MD: {sorted(missing_md)}"})

    # Verify total_pages matches
    if index_path.exists() and not any(e["message"].startswith("index.json") for e in errors):
        with open(index_path) as f:
            index = json.load(f)
        total = index.get("total_pages", 0)
        actual = max(len(png_files), len(md_files))
        if total != actual:
            errors.append({
                "severity": "warning",
                "message": f"index.json 声明 {total} 页，实际找到 {actual} 个文件"
            })

    return errors


def validate_step2(output_dir: str) -> list[dict]:
    """Validate outline generation output."""
    errors = []
    outline_path = Path(output_dir) / "02_outline" / "outline.json"

    if not outline_path.exists():
        return [{"severity": "error", "message": "02_outline/outline.json 不存在"}]

    try:
        with open(outline_path) as f:
            outline = json.load(f)
    except json.JSONDecodeError as e:
        return [{"severity": "error", "message": f"outline.json JSON 解析失败: {e}"}]

    # Check required structure
    sections = outline.get("outline", {}).get("sections", [])
    if not sections:
        errors.append({"severity": "error", "message": "outline.json 中无 sections"})
        return errors

    # Check page coverage
    all_pages = set()
    prev_end = 0
    for i, sec in enumerate(sections):
        start = sec.get("page_start", 0)
        end = sec.get("page_end", 0)

        if start <= 0 or end <= 0:
            errors.append({"severity": "error", "message": f"节 {sec.get('id', i)}: 页码无效 (start={start}, end={end})"})
            continue

        if start > end:
            errors.append({"severity": "error", "message": f"节 {sec.get('id', i)}: page_start ({start}) > page_end ({end})"})

        if prev_end > 0 and start != prev_end + 1:
            errors.append({
                "severity": "warning",
                "message": f"节 {sec.get('id', i)}: 页码不连续，上一节结束于 {prev_end}，本节从 {start} 开始"
            })

        if not sec.get("style_hint"):
            errors.append({"severity": "warning", "message": f"节 {sec.get('id', i)}: 缺少 style_hint"})

        for p in range(start, end + 1):
            all_pages.add(p)
        prev_end = end

    # Cross-check with materials
    index_path = Path(output_dir) / "01_materials" / "index.json"
    if index_path.exists():
        with open(index_path) as f:
            index = json.load(f)
        total = index.get("total_pages", 0)
        expected = set(range(1, total + 1))
        missing = expected - all_pages
        extra = all_pages - expected
        if missing:
            errors.append({"severity": "error", "message": f"大纲中缺少页面: {sorted(missing)}"})
        if extra:
            errors.append({"severity": "warning", "message": f"大纲中包含多余页面: {sorted(extra)}"})

    return errors


def validate_step3(output_dir: str) -> list[dict]:
    """Validate manuscript generation output."""
    errors = []
    manuscript_dir = Path(output_dir) / "03_manuscript"
    outline_path = Path(output_dir) / "02_outline" / "outline.json"

    if not manuscript_dir.exists():
        return [{"severity": "error", "message": "03_manuscript/ 目录不存在"}]

    # Load outline to know expected sections
    if outline_path.exists():
        with open(outline_path) as f:
            outline = json.load(f)
        sections = outline.get("outline", {}).get("sections", [])
    else:
        errors.append({"severity": "warning", "message": "无法加载 outline.json 进行交叉验证"})
        sections = []

    # Check manuscript.json
    manifest_path = manuscript_dir / "manuscript.json"
    if not manifest_path.exists():
        errors.append({"severity": "warning", "message": "manuscript.json 不存在"})

    # Check section files
    section_files = sorted(manuscript_dir.glob("section-*.md"))
    if not section_files:
        errors.append({"severity": "error", "message": "未找到任何 section-XX.md 文件"})
        return errors

    if sections and len(section_files) != len(sections):
        errors.append({
            "severity": "warning",
            "message": f"section 文件数 ({len(section_files)}) ≠ outline 节数 ({len(sections)})"
        })

    # Check PAGE markers and control instruction leaks in each section
    # These patterns catch stage directions that should never appear in TTS-ready text
    control_pattern = re.compile(
        r"（[语速停顿恢复重点强调轻声加快放慢深呼吸切换语气激动平缓严肃庆祝兴奋悲伤感慨愤怒]{2,}）"  # （语速放慢） etc
        r"|\[语速停顿恢复重点强调轻声加快放慢切换语气]{2,}\]"  # [语速加快] etc
        r"|\(语速停顿恢复重点强调轻声加快放慢切换语气]{2,}\)"  # (停顿) etc
    )
    for sec_file in section_files:
        content = sec_file.read_text(encoding="utf-8")
        page_markers = re.findall(r"\[PAGE:(\d+)\]", content)

        if not page_markers:
            errors.append({"severity": "error", "message": f"{sec_file.name}: 未找到 [PAGE:N] 标记"})
            continue

        # Extract section index
        sec_match = re.search(r"section-(\d+)", sec_file.name)
        if sec_match and sections:
            sec_idx = int(sec_match.group(1)) - 1
            if sec_idx < len(sections):
                expected_start = sections[sec_idx].get("page_start", 0)
                expected_end = sections[sec_idx].get("page_end", 0)
                actual_pages = sorted(int(p) for p in page_markers)
                expected_pages = list(range(expected_start, expected_end + 1))

                if actual_pages != expected_pages:
                    errors.append({
                        "severity": "warning",
                        "message": f"{sec_file.name}: PAGE 标记页码 {actual_pages} ≠ 预期 {expected_pages}"
                    })

        # Check for control instruction leaks
        control_hits = control_pattern.findall(content)
        if control_hits:
            errors.append({
                "severity": "warning",
                "message": f"{sec_file.name}: 发现疑似控制指令（非朗读内容）: {control_hits[:5]}"
            })

    return errors


def validate_step4(output_dir: str) -> list[dict]:
    """Validate final merged output (Markdown + JSON)."""
    errors = []
    final_dir = Path(output_dir) / "04_final"

    if not final_dir.exists():
        return [{"severity": "error", "message": "04_final/ 目录不存在"}]

    md_path = final_dir / "output.md"
    json_path = final_dir / "output.json"

    # --- Check output.md ---
    if not md_path.exists():
        errors.append({"severity": "error", "message": "output.md 不存在"})
    else:
        md_content = md_path.read_text(encoding="utf-8")
        if len(md_content.strip()) < 100:
            errors.append({"severity": "error", "message": "output.md 内容过短（< 100 字符）"})

        md_pages = sorted(set(int(m) for m in re.findall(r"\[PAGE:(\d+)\]", md_content)))
        if not md_pages:
            errors.append({"severity": "error", "message": "output.md 中未找到 [PAGE:N] 标记"})

    # --- Check output.json ---
    if not json_path.exists():
        errors.append({"severity": "error", "message": "output.json 不存在"})
    else:
        try:
            with open(json_path, encoding="utf-8") as f:
                data = json.load(f)
        except json.JSONDecodeError as e:
            errors.append({"severity": "error", "message": f"output.json JSON 解析失败: {e}"})
            data = None

        if data:
            # Check metadata
            meta = data.get("metadata", {})
            for field in ["title", "total_pages", "total_sections", "total_words"]:
                if field not in meta:
                    errors.append({"severity": "warning", "message": f"output.json metadata 缺少字段: {field}"})

            # Check sections structure
            sections = data.get("sections", [])
            if not sections:
                errors.append({"severity": "error", "message": "output.json 中无 sections"})
            else:
                json_pages = []
                for sec in sections:
                    pages = sec.get("pages", [])
                    if not pages:
                        errors.append({"severity": "warning", "message": f"output.json 节 {sec.get('id', '?')}: pages 为空"})
                    for p in pages:
                        if "page" not in p:
                            errors.append({"severity": "error", "message": f"output.json: page 条目缺少 page 字段"})
                        else:
                            json_pages.append(p["page"])
                        if not p.get("content", "").strip():
                            errors.append({"severity": "warning", "message": f"output.json: 第 {p.get('page', '?')} 页 content 为空"})

                json_pages = sorted(set(json_pages))

                # Cross-check page counts between MD and JSON
                if md_path.exists() and md_pages and json_pages:
                    if md_pages != json_pages:
                        errors.append({
                            "severity": "error",
                            "message": f"MD 与 JSON 页面不一致: MD={md_pages}, JSON={json_pages}"
                        })

    # Cross-check with materials for total page coverage
    index_path = Path(output_dir) / "01_materials" / "index.json"
    if index_path.exists():
        with open(index_path) as f:
            index = json.load(f)
        total = index.get("total_pages", 0)
        expected = set(range(1, total + 1))

        if md_path.exists() and md_pages:
            missing = expected - set(md_pages)
            if missing:
                errors.append({"severity": "error", "message": f"output.md 中缺少页面: {sorted(missing)}"})

    return errors


def validate_step5(output_dir: str) -> list[dict]:
    """Validate PPTX notes injection output."""
    errors = []
    pptx_path = Path(output_dir) / "04_final" / "presentation.pptx"

    if not pptx_path.exists():
        return [{"severity": "error", "message": "04_final/presentation.pptx 不存在"}]

    # Try to open with python-pptx
    try:
        from pptx import Presentation as PptxPresentation
    except ImportError:
        errors.append({"severity": "warning", "message": "python-pptx 未安装，无法深度验证 PPTX"})
        return errors

    try:
        prs = PptxPresentation(str(pptx_path))
    except Exception as e:
        return [{"severity": "error", "message": f"presentation.pptx 无法打开: {e}"}]

    # Load output.json to know which pages should have manuscript
    json_path = Path(output_dir) / "04_final" / "output.json"
    expected_pages = set()
    if json_path.exists():
        try:
            with open(json_path, encoding="utf-8") as f:
                data = json.load(f)
            for sec in data.get("sections", []):
                for page in sec.get("pages", []):
                    if page.get("content", "").strip():
                        expected_pages.add(page["page"])
        except Exception:
            errors.append({"severity": "warning", "message": "无法读取 output.json 进行交叉验证"})

    # Check notes in PPTX
    manuscript_count = 0
    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
        except Exception:
            notes_text = ""

        if "<manuscript>" in notes_text:
            manuscript_count += 1

            # Check it also has closing tag
            if "</manuscript>" not in notes_text:
                errors.append({"severity": "warning", "message": f"第 {page_num} 页: 有 <manuscript> 但缺少 </manuscript>"})
        elif page_num in expected_pages:
            errors.append({"severity": "warning", "message": f"第 {page_num} 页: output.json 中有内容但 PPTX 备注中无 <manuscript>"})

    if expected_pages and manuscript_count == 0:
        errors.append({"severity": "error", "message": "PPTX 中未找到任何 <manuscript> 标记"})
    elif expected_pages and manuscript_count != len(expected_pages):
        errors.append({
            "severity": "warning",
            "message": f"<manuscript> 数量 ({manuscript_count}) ≠ output.json 页面数 ({len(expected_pages)})"
        })

    return errors


def main():
    parser = argparse.ArgumentParser(description="验证 PPT-TTS 流水线各步骤输出")
    parser.add_argument("output_dir", help="工作目录路径")
    parser.add_argument("--step", type=int, required=True, choices=[1, 2, 3, 4, 5], help="验证哪一步的输出")
    args = parser.parse_args()

    validators = {
        1: validate_step1,
        2: validate_step2,
        3: validate_step3,
        4: validate_step4,
        5: validate_step5,
    }

    step_names = {
        1: "素材提取",
        2: "大纲生成",
        3: "逐字稿生成",
        4: "合并输出",
        5: "备注回写",
    }

    errors = validators[args.step](args.output_dir)

    print(f"\n{'='*50}")
    print(f"步骤 {args.step} ({step_names[args.step]}) 验证结果")
    print(f"{'='*50}")

    if not errors:
        print("\n✅ 验证通过，无问题。")
        print("\n--- JSON ---")
        print(json.dumps({"step": args.step, "status": "pass", "errors": []}, ensure_ascii=False))
        sys.exit(0)

    err_count = sum(1 for e in errors if e["severity"] == "error")
    warn_count = sum(1 for e in errors if e["severity"] == "warning")

    for e in errors:
        icon = "❌" if e["severity"] == "error" else "⚠️"
        print(f"\n{icon} [{e['severity'].upper()}] {e['message']}")

    print(f"\n总结: {err_count} 个错误, {warn_count} 个警告")

    print("\n--- JSON ---")
    print(json.dumps({
        "step": args.step,
        "status": "fail" if err_count > 0 else "warn",
        "errors": errors,
    }, ensure_ascii=False, indent=2))

    sys.exit(1 if err_count > 0 else 0)


if __name__ == "__main__":
    main()
