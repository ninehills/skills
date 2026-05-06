#!/usr/bin/env python3
"""Extract PPT slide content and speaker notes to markdown files."""

import sys
import os
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def extract_slide_content(slide, slide_num):
    """Extract slide content and speaker notes."""
    lines = []

    # Title
    title = slide.shapes.title
    if title and title.text.strip():
        lines.append(f"# {title.text.strip()}")

    # Body content - extract all text from shapes
    for shape in slide.shapes:
        # Skip title (already handled)
        if shape == title:
            continue

        # Skip notes slide placeholder
        try:
            if shape.placeholder_format.type == 1:  # notes placeholder
                continue
        except (ValueError, AttributeError):
            pass

        # Extract text from shape
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text:
                lines.append(text)

    # Speaker notes
    notes_slide = slide.notes_slide
    if notes_slide:
        notes_text = notes_slide.notes_text_frame.text.strip()
        if notes_text:
            lines.append(f"\n**演讲者备注:**\n{notes_text}")

    return "\n\n".join(lines)


def check_dependencies():
    """Check if required dependencies are installed."""
    errors = []

    # Check python-pptx
    try:
        from pptx import Presentation
    except ImportError:
        errors.append("python-pptx: pip install python-pptx")

    # Check LibreOffice
    if os.system("which soffice > /dev/null 2>&1") != 0:
        errors.append("LibreOffice (soffice): brew install --cask libreoffice")

    # Check poppler
    if os.system("which pdftoppm > /dev/null 2>&1") != 0:
        errors.append("poppler (pdftoppm): brew install poppler")

    return errors


def main():
    if len(sys.argv) < 3:
        print("Usage: extract_ppt.py <pptx_path> <output_dir>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2]

    # Check dependencies
    errors = check_dependencies()
    if errors:
        print("ERROR: Missing dependencies:")
        for err in errors:
            print(f"  - {err}")
        sys.exit(1)

    # Create output directory
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    # Extract text to markdown
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        content = extract_slide_content(slide, i + 1)
        md_path = os.path.join(output_dir, f"slide-{i + 1:02d}.md")
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(content)

    print(f"Extracted {len(prs.slides)} slides to {output_dir}")


if __name__ == "__main__":
    main()
